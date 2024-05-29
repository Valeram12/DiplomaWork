import aspose.slides as slides
import os
import g4f
from pptx import Presentation
def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def presentationAnalyze(namePres):
    print(namePres)
    is_ppt = False
    if '.pptx' not in namePres:
        is_ppt = True
        pres = slides.Presentation(namePres)
        namePres = namePres.replace('.ppt', '.pptx')
        pres.save(namePres, slides.export.SaveFormat.PPTX)
    print("повинна початись отримання тексту")
    pptx_text = extract_text_from_pptx(namePres)
    cleaned_text = pptx_text.replace("Evaluation only.", "") \
                            .replace("Created with Aspose.Slides for Python via .NET 24.4.", "") \
                            .replace("Copyright 2004-2024Aspose Pty Ltd.", "") \
                            .replace("\n", "")
    prompt = (
            "Ти повенен проаналізувати текст та дати дуже короткий опис цього тексту головну тему та основні пункти. "
            "твоя відповідь повинна бути максимально коротка та структурно розписана без самостійного опису цих тем і "
            "пунктів, без будь-якої стилізації текста, звичайний розмір та шрифт" + cleaned_text)
    response = g4f.ChatCompletion.create(
        model=g4f.models.gpt_4,
        messages=[{"role": "user", "content": prompt}],
    )

    result = response.replace("*", '')
    if is_ppt:
        os.remove(namePres)
    return result

# print(presentationAnalyze("pptx/МедведчукВалерійІПС-41TakamisawaCybernetics.pptx"))
